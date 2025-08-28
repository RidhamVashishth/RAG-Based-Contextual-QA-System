import os
from dotenv import load_dotenv
import streamlit as st
import google.generativeai as genai
from pypdf import PdfReader
import docx
import pptx
import openpyxl
from langchain.schema import Document
from langchain_text_splitters import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings


# ------------------------🔐 API CONFIGURATION ------------------------
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-2.0-flash")

# ------------------------🧠 EMBEDDING MODEL ------------------------
@st.cache_resource(show_spinner="🔄 Loading the HuggingFace Embedding Model...")
def get_embedding_model():
    """Loads the HuggingFace embedding model."""
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

def extract_text_from_file(uploaded_file):
    """Extracts raw text from various file types."""
    raw_text = ""
    # Get the file extension from the uploaded file's name
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()

    if file_extension == ".pdf":
        pdf = PdfReader(uploaded_file)
        for page in pdf.pages:
            raw_text += page.extract_text() or ""
    elif file_extension == ".docx":
        doc = docx.Document(uploaded_file)
        for para in doc.paragraphs:
            raw_text += para.text + "\n"
    elif file_extension == ".pptx":
        prs = pptx.Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text += shape.text + "\n"
    elif file_extension == ".xlsx":
        workbook = openpyxl.load_workbook(uploaded_file)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                # Join cell values, converting None to an empty string for clean text
                raw_text += " ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
    
    return raw_text

# ------------------------🎯 STREAMLIT FRONT-END ------------------------
st.title("📄 AskTheDoc - Multi-Format Q&A System")
st.markdown("### 🤖 Powered by HuggingFace Embeddings + FAISS + Gemini")
st.markdown("---")

# File Uploader for multiple types
uploaded_file = st.file_uploader(
    "📤 Upload your PDF, DOCX, PPTX, or XLSX Document",
    type=["pdf", "docx", "pptx", "xlsx"]
)

# ------------------------📄 DOCUMENT PROCESSING ------------------------
if uploaded_file:
    # Use the universal text extraction function
    raw_text = extract_text_from_file(uploaded_file)

    if not raw_text.strip():
        st.error("Could not extract text from the document. The file might be empty or corrupted.")
    else:
        # Wrap the extracted text in a Document object
        document = Document(page_content=raw_text)

        # Chunking the document
        splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents([document])

        if not chunks:
            st.error("Failed to split the document into chunks.")
        else:
            # Prepare FAISS vector DB
            texts = [chunk.page_content for chunk in chunks]
            embedding_model = get_embedding_model()
            vector_db = FAISS.from_texts(texts, embedding_model)
            retriever = vector_db.as_retriever()

            # User Input Form
            st.success("✅ Document processed successfully. Ask your questions below!")

            with st.form(key="query_form", clear_on_submit=True):
                user_input = st.text_input("💬 Enter your question:", key="user_input")
                submitted = st.form_submit_button("🔍 Submit")

            if submitted and user_input:
                with st.chat_message("user"):
                    st.write(user_input)

                with st.spinner("🧠 Analyzing and thinking..."):
                    # Retrieve relevant chunks from the vector database
                    retrieved_docs = retriever.get_relevant_documents(user_input)
                    context = "\n\n".join(doc.page_content for doc in retrieved_docs)

                    # Prompt formatting
                    prompt = f"""You are an expert research assistant.Your ONLY source of information is the context provided.
RULE: You MUST NOT access external knowledge.
RULE: If the context does not contain the answer, you MUST reply with 'The provided text does not contain this information.'

Context:
---
{context}
---

Query: {user_input}

Answer:"""

                    # Generate the answer using the Gemini model
                    response = model.generate_content(prompt)

                    # Show response
                    with st.chat_message("assistant"):
                        st.markdown(response.text)
else:
    st.info("⚠️ Please upload a document to begin.")
